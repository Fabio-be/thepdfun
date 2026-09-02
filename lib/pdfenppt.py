import fitz
from pptx import Presentation
import tempfile
import os

def convert_toppt(input_path, output_path):
    pdf = fitz.open(input_path)
    prs = Presentation()

    with tempfile.TemporaryDirectory() as temp_dir:
        for i, page in enumerate(pdf):
            img_path = os.path.join(temp_dir, f"page_{i}.png")

            pix = page.get_pixmap()
            pix.save(img_path)

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            slide.shapes.add_picture(
                img_path,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height
            )

    prs.save(output_path)
